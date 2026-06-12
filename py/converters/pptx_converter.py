from pptx import Presentation
import re
from converters import content_detector
from converters import integrity

_CONTENT_TYPE_MAP = {
    "financial": "Financial Statement",
    "medical":   "Medical Document",
    "legal":     "Legal Document",
    "real_estate": "Real Estate Document",
    "insurance": "Insurance Document",
}

_FNAME_KEYWORDS = {
    "Insurance Document": ['insurance', 'policy', 'coverage', 'premium'],
    "Medical Document":   ['medical', 'health', 'lab', 'clinical', 'hospital'],
    "Legal Document":     ['contract', 'agreement', 'legal', 'brief', 'motion'],
    "Financial Statement":['financial', 'budget', 'revenue', 'forecast', 'statement'],
}


def _detect_content_type(full_text: str, filename: str) -> str:
    detected = content_detector.detect_content_type(full_text)
    ct = _CONTENT_TYPE_MAP.get(detected, "PowerPoint Presentation")
    if ct == "PowerPoint Presentation":
        fname_lower = filename.lower()
        for content_type, keywords in _FNAME_KEYWORDS.items():
            if any(k in fname_lower for k in keywords):
                return content_type
    return ct


def convert(file_path: str, filename: str, process_date: str, file_slug: str = "") -> tuple[str, dict]:
    prs = Presentation(file_path)

    title = filename
    if prs.slides:
        first_slide = prs.slides[0]
        for shape in first_slide.shapes:
            if shape.has_text_frame and shape.shape_type == 13:
                title = shape.text_frame.text.strip() or filename
                break
            # shape.placeholder_format raises ValueError on non-placeholder shapes,
            # so gate on is_placeholder first (common in real/converted decks).
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                title = shape.text_frame.text.strip() or filename
                break

    lines = []
    full_text_parts = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_title = ""
        slide_body = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            full_text_parts.append(text)
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                slide_title = text
                continue
            slide_body.append(text)

        anchor_id = f"{file_slug}-page-{i:03d}" if file_slug else f"page-{i:03d}"
        lines.append(f'\n<a id="{anchor_id}"></a>\n\n### Slide {i}' + (f": {slide_title}" if slide_title else ""))
        lines.append("")
        for block in slide_body:
            for line in block.splitlines():
                stripped = line.strip()
                if stripped:
                    lines.append(f"- {stripped}")
        lines.append("")

    full_text = "\n".join(full_text_parts)
    content_type = _detect_content_type(full_text, filename)

    years = re.findall(r'\b(20\d{2}|19\d{2})\b', full_text[:2000])

    meta = {
        "potential_titles": [title] if title != filename else [],
        "form": "",
        "form_name": "",
        "tax_year": "",
        "case_numbers": [],
        "page_count": len(prs.slides),
        "content_type": content_type,
        "scanned_ratio": 0,
        "ocr_pages": 0,
        "ocr_available": False,
    }

    combined = "\n".join(lines)
    meta["integrity"] = integrity.check_numbers(full_text, combined)
    return combined, meta
